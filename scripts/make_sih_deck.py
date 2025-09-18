import argparse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor

DEFAULT_TITLE = "Find My Stipend — SIH Pitch"
DEFAULT_TEAM = "Team Name"
DEFAULT_COLLEGE = "College Name"
DEFAULT_CONTACT = "Email • Phone • GitHub: Ruchin-Audichya/StudentPilot"
LIVE_FRONTEND = "https://wms-virid-six.vercel.app"
BACKEND_HEALTH = "https://studentpilot.onrender.com/health"

SLIDES = [
    {
        "title": "Find My Stipend — AI Career Copilot",
        "bullets": [
            "Smart India Hackathon — Software + AI",
            f"Live Frontend: {LIVE_FRONTEND}",
            f"Backend Health: {BACKEND_HEALTH}",
        ],
        "is_title": True,
    },
    {
        "title": "Problem & Why Now",
        "bullets": [
            "Students struggle to find verified, relevant internships quickly",
            "Fragmented sources: Internshala, LinkedIn, company ATS, government portals",
            "Low-quality matches waste time; ATS mismatch reduces selection odds",
            "Recruiter discovery + targeted outreach are manual and slow",
        ],
    },
    {
        "title": "SIH Problem Fit",
        "bullets": [
            "AI-powered tool improving employability and access",
            "Aggregates multi-source internships + government snapshot",
            "Resume-aware ranking, HR tools, mock interviews, instant portfolio",
            "Built for Indian students: verified sources, mobile-first, low-friction",
        ],
    },
    {
        "title": "Solution Overview",
        "bullets": [
            "Multi-source aggregation (LinkedIn, Internshala, ATS, Govt)",
            "Resume-aware ranking & gap analysis (skills vs JD)",
            "HR tools: recruiter links & profiles for targeted outreach",
            "Mock interviews with voice + AI feedback",
            "Instant portfolio ZIP (Vercel-ready)",
        ],
    },
    {
        "title": "Key Differentiators",
        "bullets": [
            "Resume-aware scoring → saves time, higher relevancy",
            "Verified gov snapshot (AICTE/NCS/MyGov/DRDO/NITI/Maha)",
            "HR profiles + search links → better response rates",
            "Zero-config portfolio deploy (vercel.json included)",
            "Mobile-first UX with safe-area handling",
        ],
    },
    {
        "title": "Architecture (High-level)",
        "bullets": [
            "Frontend: React + Vite + Tailwind (Vercel)",
            "Backend: FastAPI (Python), scraping + enrichment (Render/EB)",
            "Sources: Internshala, LinkedIn, ATS (Lever/Greenhouse/Workday), Gov feeds",
            "LLM: Gemini/OpenRouter (pluggable) for analyzer, messages, portfolio",
            "Vercel serverless proxy → backend origin (secure cross-origin)",
        ],
    },
    {
        "title": "Live Demo Flow",
        "bullets": [
            "Upload resume → extract skills/roles",
            "Search → ranked internships with ATS/Gov labels",
            "Analyzer → missing keywords + ATS score",
            "HR tools → recruiter links/profiles for a top company",
            "Mock Interview → voice Q&A + feedback",
            "Generate portfolio ZIP → deploy on Vercel",
        ],
    },
    {
        "title": "Government Snapshot (SIH Lens)",
        "bullets": [
            "Aggregates AICTE/NCS/MyGov/DRDO/NITI/MahaSwayam",
            "Filters by state, verified-only; export CSV/JSON",
            "Cache + refresh; admin moderation endpoint",
            "Authentic, accessible opportunities for Tier-2/3 students",
        ],
    },
    {
        "title": "Resume Genius (ATS Optimizer)",
        "bullets": [
            "Parses resume + JD → ATS score (0–100)",
            "Matched/missing skills; actionable suggestions",
            "Project ideas for skill gaps; tailored resume outline",
            "Optional messages: cover letter / LinkedIn note",
        ],
    },
    {
        "title": "HR Tools & Outreach",
        "bullets": [
            "Smart recruiter search links (site:linkedin patterns)",
            "HR profiles for ATS companies from results",
            "Shift from mass-apply → targeted outreach",
            "Expected uplift: +15–25% response vs generic cold",
        ],
    },
    {
        "title": "Impact & KPIs (3–6 months)",
        "bullets": [
            "Time-to-first-good-match: 2–3 weeks → < 3 days",
            "Application-to-response rate: 1–2% → 6–10%",
            "ATS score improvement: +20–30 points",
            "Portfolio deploy rate: 60% of active users",
            "Gov verified listings: 500–2000/month aggregated",
        ],
    },
    {
        "title": "Risks & Mitigation",
        "bullets": [
            "Data limits → fallback to gov feeds, ATS sources",
            "Scraping constraints → throttle + link-based patterns",
            "LLM variability → deterministic parsing + timeouts",
            "Security/PII → proxy via serverless, scoped CORS, anonymized logs",
            "Scalability → stateless FE, horizontal BE scale, caching",
        ],
    },
    {
        "title": "Execution Timeline",
        "bullets": [
            "W1–2: Gov snapshot + analyzer hardening",
            "W3–4: HR tools maturity; mobile polish",
            "W5: Mock Interview latency + caching",
            "W6: College pilot (100–300 users); measure KPIs; finalize docs",
            "Next: Placement cell dashboards; employer verification",
        ],
    },
    {
        "title": "Call to Action",
        "bullets": [
            "Working product aligned with SIH outcomes",
            "Ready to pilot with placement cells; deploy < 1 day",
            "Ask: Support for campus integrations and scaling credits",
            f"Live: {LIVE_FRONTEND}",
        ],
    },
]


def add_title_slide(prs: Presentation, title_text: str, subtitle_text: str):
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = title_text
    subtitle.text = subtitle_text


def add_bullets_slide(prs: Presentation, title_text: str, bullets: list[str]):
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]

    title.text = title_text

    tf = body.text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(20)

    # Small visual polish: left alignment and color
    for p in tf.paragraphs:
        p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
        p.font.color.rgb = RGBColor(0x22, 0x22, 0x22)


def build_presentation(team: str, college: str, contact: str, out_path: str):
    prs = Presentation()

    # Title slide
    title_text = DEFAULT_TITLE
    subtitle_text = f"{team} • {college}\n{contact}"
    add_title_slide(prs, title_text, subtitle_text)

    # Content slides
    for s in SLIDES:
        if s.get("is_title"):
            # Convert the first slide after title slide into normal content slide
            add_bullets_slide(prs, s["title"], s["bullets"])  # Already title slide above
        else:
            add_bullets_slide(prs, s["title"], s["bullets"])  

    prs.save(out_path)


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Generate SIH PPT deck for Find My Stipend")
    parser.add_argument("--team", default=DEFAULT_TEAM)
    parser.add_argument("--college", default=DEFAULT_COLLEGE)
    parser.add_argument("--contact", default=DEFAULT_CONTACT)
    parser.add_argument("--out", default="SIH_FindMyStipend.pptx")
    args = parser.parse_args()

    build_presentation(args.team, args.college, args.contact, args.out)
    print(f"Deck generated: {args.out}")
